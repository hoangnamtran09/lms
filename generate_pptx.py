#!/usr/bin/env python3
"""Generate EduSelf presentation PPTX from HTML content."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ── Constants ──
BG_DARK   = RGBColor(0x06, 0x06, 0x12)
BG_CARD   = RGBColor(0x0F, 0x0F, 0x24)
TEXT_WHITE = RGBColor(0xF8, 0xFA, 0xFC)
TEXT_GRAY  = RGBColor(0x94, 0xA3, 0xB8)
TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)
BLUE       = RGBColor(0x60, 0xA5, 0xFA)
PURPLE     = RGBColor(0xA7, 0x8B, 0xFA)
GREEN      = RGBColor(0x34, 0xD3, 0x99)
AMBER      = RGBColor(0xFB, 0xBF, 0x24)
ROSE       = RGBColor(0xFB, 0x71, 0x85)
RED_LIGHT  = RGBColor(0xFC, 0xA5, 0xA5)
GREEN_LIGHT = RGBColor(0x6E, 0xE7, 0xB7)
WHITE_10   = RGBColor(0x1A, 0x1A, 0x2E)

W = Inches(13.333)  # 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

# ── Helpers ──

def slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=TEXT_WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='Be Vietnam Pro'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox, tf

def add_para(tf, text, font_size=14, color=TEXT_WHITE, bold=False, align=PP_ALIGN.LEFT,
             space_before=Pt(4), space_after=Pt(2)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Be Vietnam Pro'
    p.alignment = align
    p.space_before = space_before
    p.space_after = space_after
    return p

def add_rect(slide, left, top, width, height, fill_color=BG_CARD, border_color=None,
             corner_radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_circle(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_slide_number(slide, num, total=27):
    add_text_box(slide, 12.3, 7.05, 0.8, 0.35, f'{num}/{total}',
                 font_size=9, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)

def add_section_bg(slide, num):
    slide_bg(slide)
    # Large number watermark
    add_text_box(slide, 0, 1.5, 13.333, 3.0, str(num).zfill(2),
                 font_size=140, color=RGBColor(0x0D, 0x15, 0x30), bold=True, align=PP_ALIGN.CENTER)

def set_transition(slide, type='fade'):
    """Set slide transition. type: fade, push, morph, wipe, etc."""
    mapping = {
        'fade': 'fade',
    }
    # python-pptx limited transition support, using basic
    slide.slide_layout = prs.slide_layouts[6]  # blank

# ═══════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
# Logo mark placeholder
r = add_rect(sl, 5.8, 1.3, 1.6, 1.6, fill_color=RGBColor(0x0D, 0x18, 0x3A),
             border_color=RGBColor(0x25, 0x25, 0x40), corner_radius=0.3)
add_text_box(sl, 5.8, 1.55, 1.6, 0.8, '🎓', font_size=44, align=PP_ALIGN.CENTER)
# Title
add_text_box(sl, 1.5, 3.3, 10.3, 2.2, 'Hệ thống Quản lý\nHọc tập Thông minh',
             font_size=54, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(sl, 2.5, 5.2, 8.3, 0.8,
             'Nền tảng học tập cá nhân hoá ứng dụng Trí tuệ Nhân tạo dành cho học sinh Việt Nam',
             font_size=18, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
# Tags
for i, (text, clr) in enumerate([
    ('🤖 Trợ lý AI thông minh', RGBColor(0xC4, 0xB5, 0xFD)),
    ('📚 Quản lý học tập toàn diện', RGBColor(0x93, 0xC5, 0xFD)),
    ('🏆 Học mà chơi & Thành tựu', RGBColor(0x6E, 0xE7, 0xB7)),
]):
    add_text_box(sl, 3.2 + i * 3.0, 6.3, 2.8, 0.4, text, font_size=12, color=clr, align=PP_ALIGN.CENTER)
add_slide_number(sl, 1)

# ═══════════════════════════════════════════════════════════
# SLIDE 2: SECTION — Giới thiệu
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_section_bg(sl, 1)
add_text_box(sl, 0, 2.8, 13.333, 1.0, 'Giới thiệu về EduSelf',
             font_size=44, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(sl, 2.0, 4.0, 9.3, 1.5,
             'EduSelf là nền tảng học tập trực tuyến thông minh, kết hợp công nghệ AI hiện đại với phương pháp sư phạm gợi mở, giúp cá nhân hoá trải nghiệm học tập cho từng học sinh Việt Nam.',
             font_size=18, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
add_slide_number(sl, 2)

# ═══════════════════════════════════════════════════════════
# SLIDE 3: TỔNG QUAN NỀN TẢNG
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
add_text_box(sl, 0.5, 0.3, 12.3, 0.8, 'EduSelf là nền tảng toàn diện cho giáo dục',
             font_size=34, bold=True, align=PP_ALIGN.CENTER)
# Stats
stats = [('5', 'Vai trò người dùng'), ('45+', 'API Endpoints'), ('10', 'Nhóm tính năng'), ('24/7', 'AI sẵn sàng')]
for i, (num, label) in enumerate(stats):
    x = 1.5 + i * 2.8
    add_text_box(sl, x, 1.3, 2.5, 0.8, num, font_size=44, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(sl, x, 2.0, 2.5, 0.4, label, font_size=14, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
# Role cards
roles = [
    ('👑', 'Super Admin', 'Quản trị hệ thống cao nhất,\ntoàn quyền kiểm soát'),
    ('⚙️', 'Admin', 'Quản lý người dùng, môn học,\nbài tập & xuất dữ liệu'),
    ('👩‍🏫', 'Giáo viên', 'Tạo bài tập, chấm điểm,\ntheo dõi học sinh, điểm danh'),
    ('👨‍👩‍👧', 'Phụ huynh', 'Theo dõi tiến độ học tập\ncủa con, nhận báo cáo AI'),
    ('🎓', 'Học sinh', 'Học bài, chat AI, làm quiz,\nnộp bài tập, tích luỹ thành tựu'),
]
for i, (icon, name, desc) in enumerate(roles):
    x = 0.6 + i * 2.55
    add_rect(sl, x, 2.8, 2.35, 3.2, corner_radius=0.15)
    add_text_box(sl, x, 3.0, 2.35, 0.5, icon, font_size=28, align=PP_ALIGN.CENTER)
    add_text_box(sl, x, 3.6, 2.35, 0.4, name, font_size=14, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(sl, x, 4.1, 2.35, 0.9, desc, font_size=10, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
add_slide_number(sl, 3)

# ═══════════════════════════════════════════════════════════
# SLIDE 4: SECTION — Lý do chọn đề tài
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_section_bg(sl, 2)
add_text_box(sl, 0, 2.8, 13.333, 1.0, 'Lý do chọn đề tài',
             font_size=44, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(sl, 1.5, 4.0, 10.3, 1.5,
             'Từ thực tiễn giáo dục Việt Nam đến khát vọng tạo ra sản phẩm công nghệ "Make in Vietnam" phục vụ học sinh, giáo viên và phụ huynh',
             font_size=18, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
add_slide_number(sl, 4)

# ═══════════════════════════════════════════════════════════
# SLIDE 5: LÝ DO CHI TIẾT
# ═══════════════════════════════════════════════════════════
reasons = [
    ('📱', 'Nhu cầu chuyển đổi số giáo dục', 'Việt Nam đẩy mạnh chuyển đổi số\ntheo đề án Chính phủ. Cần nền tảng\nhọc tập số thực sự hiệu quả.'),
    ('📍', 'Khoảng cách thành thị – nông thôn', 'Học sinh vùng sâu thiếu giáo viên\ngiỏi & tài liệu chất lượng.\nAI có thể thu hẹp khoảng cách này.'),
    ('🤖', 'Tiềm năng đột phá của AI', 'Trí tuệ nhân tạo đã đủ trưởng thành\nđể làm gia sư ảo. gợi mở + AI =\ntrải nghiệm học tập chưa từng có.'),
    ('📋', 'Hạn chế của LMS hiện có', 'Moodle, Google Classroom... thiếu AI,\nthiếu gamification, giao diện phức tạp,\nkhông Việt hoá đầy đủ.'),
    ('📊', 'Thói quen học tập sau COVID-19', 'Học trực tuyến là tất yếu. Học sinh cần\nhơn cả video call — cần nền tảng\ntương tác, thông minh, tạo động lực.'),
    ('🇻🇳', 'Khát vọng "Make in Vietnam"', 'Chứng minh người Việt làm được sản\nphẩm công nghệ giáo dục đẳng cấp,\nbằng tiếng Việt, cho người Việt.'),
]
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
add_text_box(sl, 0.3, 0.15, 12.7, 0.7, 'Tại sao chúng tôi chọn xây dựng EduSelf?',
             font_size=30, bold=True, align=PP_ALIGN.CENTER)
for i, (icon, title, desc) in enumerate(reasons):
    row = i // 3
    col = i % 3
    x = 0.5 + col * 4.2
    y = 1.1 + row * 3.1
    add_rect(sl, x, y, 3.9, 2.8, corner_radius=0.15)
    add_text_box(sl, x + 0.2, y + 0.2, 0.5, 0.5, icon, font_size=22)
    add_text_box(sl, x + 0.8, y + 0.2, 2.8, 0.5, title, font_size=14, color=TEXT_WHITE, bold=True)
    add_text_box(sl, x + 0.2, y + 1.0, 3.5, 1.6, desc, font_size=10.5, color=TEXT_GRAY)
add_slide_number(sl, 5)

# ═══════════════════════════════════════════════════════════
# SLIDE 6: SECTION — Vấn đề
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_section_bg(sl, 7)
add_text_box(sl, 0, 2.8, 13.333, 1.0, 'Thực trạng & Vấn đề',
             font_size=44, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(sl, 2.0, 4.0, 9.3, 1.5,
             'Những khó khăn trong giáo dục truyền thống mà học sinh, giáo viên và phụ huynh đang đối mặt mỗi ngày',
             font_size=18, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
add_slide_number(sl, 4)

# ═══════════════════════════════════════════════════════════
# SLIDE 5: VẤN ĐỀ CHI TIẾT
# ═══════════════════════════════════════════════════════════
problems = [
    ('📖', 'Học tập thụ động, một chiều', 'Học sinh chỉ đọc và ghi nhớ. Không có người hướng dẫn khi tự học ở nhà. Thắc mắc không được giải đáp kịp thời.'),
    ('👥', 'Thiếu cá nhân hoá', 'Một giáo viên dạy hàng chục học sinh cùng một phương pháp. Học sinh giỏi nhàm chán, học sinh yếu không theo kịp.'),
    ('📝', 'Chấm điểm thủ công, nặng nhọc', 'Giáo viên mất hàng giờ soạn đề, chấm bài, vào điểm. Thời gian quý báu cho giảng dạy bị cắt giảm.'),
    ('📉', 'Không theo dõi được tiến độ', 'Phụ huynh không biết con học gì, học bao lâu. Nhà trường không có dữ liệu để đánh giá chất lượng.'),
    ('🎯', 'Không phát hiện lỗ hổng kiến thức', 'Học sinh không biết mình yếu phần nào để ôn tập. Lỗ hổng tích luỹ âm thầm qua thời gian.'),
    ('😴', 'Thiếu động lực học tập', 'Không có yếu tố khuyến khích hay khen thưởng. Học sinh dễ chán nản, thiếu cạnh tranh lành mạnh.'),
]
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
add_text_box(sl, 0.5, 0.2, 12.3, 0.7, 'Những vấn đề trước khi có EduSelf',
             font_size=32, bold=True, align=PP_ALIGN.CENTER)
for i, (icon, title, desc) in enumerate(problems):
    row = i // 3
    col = i % 3
    x = 0.5 + col * 4.2
    y = 1.2 + row * 3.1
    add_rect(sl, x, y, 3.9, 2.8, corner_radius=0.15)
    add_text_box(sl, x + 0.2, y + 0.2, 0.5, 0.5, icon, font_size=22)
    add_text_box(sl, x + 0.8, y + 0.2, 2.8, 0.5, title, font_size=15, color=TEXT_WHITE, bold=True)
    add_text_box(sl, x + 0.2, y + 1.0, 3.5, 1.6, desc, font_size=11, color=TEXT_GRAY)
add_slide_number(sl, 5)

# ═══════════════════════════════════════════════════════════
# SLIDE 6: SECTION — Giải pháp
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_section_bg(sl, 7)
add_text_box(sl, 0, 2.8, 13.333, 1.0, 'Giải pháp của EduSelf',
             font_size=44, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(sl, 2.0, 4.0, 9.3, 1.5,
             'Cách EduSelf giải quyết từng vấn đề bằng công nghệ AI tiên tiến và thiết kế lấy người học làm trung tâm',
             font_size=18, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
add_slide_number(sl, 6)

# ═══════════════════════════════════════════════════════════
# SLIDE 7: GIẢI PHÁP CHI TIẾT
# ═══════════════════════════════════════════════════════════
solutions = [
    ('💬', 'Trợ lý AI thông minh — Học chủ động', 'Trợ lý AI 24/7 đặt câu hỏi gợi mở, dẫn dắt học sinh tự khám phá kiến thức.'),
    ('🧠', 'Lộ trình học cá nhân hoá', 'AI phân tích năng lực từng học sinh, đề xuất lộ trình phù hợp và bài tập riêng biệt.'),
    ('⚡', 'Tự động tạo & chấm bài', 'AI tạo câu hỏi + tự luận, tự động chấm điểm theo rubric. Giáo viên tiết kiệm 70% thời gian.'),
    ('📈', 'Theo dõi tiến độ minh bạch', 'Dashboard trực quan với biểu đồ. Phụ huynh nhận báo cáo AI hàng tuần.'),
    ('🔍', 'Chẩn đoán điểm yếu thông minh', 'Phát hiện lỗ hổng kiến thức qua bài sai, AI tạo bài tập riêng để khắc phục.'),
    ('🏆', 'Học mà chơi — Học mà chơi', 'Kim cương, streak, thành tựu và bảng xếp hạng tạo động lực học tập.'),
]
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
add_text_box(sl, 0.5, 0.2, 12.3, 0.7, 'EduSelf đã giải quyết những vấn đề đó như thế nào?',
             font_size=30, bold=True, align=PP_ALIGN.CENTER)
for i, (icon, title, desc) in enumerate(solutions):
    row = i // 2
    col = i % 2
    x = 0.5 + col * 6.3
    y = 1.2 + row * 2.0
    add_rect(sl, x, y, 5.9, 1.75, corner_radius=0.12)
    add_text_box(sl, x + 0.15, y + 0.15, 0.5, 0.5, icon, font_size=20)
    add_text_box(sl, x + 0.75, y + 0.15, 5.0, 0.4, title, font_size=14, color=TEXT_WHITE, bold=True)
    add_text_box(sl, x + 0.75, y + 0.65, 5.0, 0.9, desc, font_size=11, color=TEXT_GRAY)
add_slide_number(sl, 7)

# ═══════════════════════════════════════════════════════════
# SLIDE 8: SO SÁNH TRƯỚC/SAU
# ═══════════════════════════════════════════════════════════
before_items = [
    'Học sinh tự học một mình, không ai hướng dẫn khi gặp khó',
    'Giáo viên mất 4-5 tiếng/tuần soạn đề và chấm bài thủ công',
    'Phụ huynh mù mờ không biết con học hành ra sao',
    'Học sinh không biết mình yếu chỗ nào để cải thiện',
    'Không có động lực, dễ bỏ cuộc giữa chừng',
    'Dữ liệu học tập phân tán, không đo lường được',
]
after_items = [
    'AI Tutor đồng hành 24/7, giải đáp mọi thắc mắc',
    'AI tạo & chấm bài tự động, giáo viên chỉ kiểm tra lại',
    'Dashboard & báo cáo AI hàng tuần cho phụ huynh',
    'Chẩn đoán điểm yếu + bài tập cải thiện riêng',
    'Kim cương, streak, thành tựu, bảng xếp hạng thúc đẩy',
    'Toàn bộ dữ liệu tập trung, phân tích & xuất báo cáo',
]
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
add_text_box(sl, 1.0, 0.2, 11.3, 0.7, 'Trước và sau khi có EduSelf',
             font_size=34, bold=True, align=PP_ALIGN.CENTER)
# Before box
add_rect(sl, 0.3, 1.2, 5.8, 5.2, fill_color=RGBColor(0x1A, 0x08, 0x08),
         border_color=RGBColor(0x40, 0x10, 0x10), corner_radius=0.15)
add_text_box(sl, 0.5, 1.35, 5.4, 0.4, '✕ Trước khi có EduSelf', font_size=16, color=RED_LIGHT, bold=True)
for j, item in enumerate(before_items):
    add_text_box(sl, 0.6, 1.9 + j * 0.7, 5.2, 0.6, f'✕  {item}', font_size=11.5, color=TEXT_GRAY)
# Arrow
add_text_box(sl, 6.1, 3.2, 1.0, 0.6, '⟶', font_size=40, color=BLUE, align=PP_ALIGN.CENTER)
# After box
add_rect(sl, 7.15, 1.2, 5.85, 5.2, fill_color=RGBColor(0x08, 0x1A, 0x08),
         border_color=RGBColor(0x10, 0x40, 0x10), corner_radius=0.15)
add_text_box(sl, 7.35, 1.35, 5.4, 0.4, '✓ Sau khi có EduSelf', font_size=16, color=GREEN_LIGHT, bold=True)
for j, item in enumerate(after_items):
    add_text_box(sl, 7.45, 1.9 + j * 0.7, 5.3, 0.6, f'✓  {item}', font_size=11.5, color=GREEN_LIGHT)
add_slide_number(sl, 8)

# ═══════════════════════════════════════════════════════════
# SLIDE 9: SECTION — Tính năng
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_section_bg(sl, 7)
add_text_box(sl, 0, 2.8, 13.333, 1.0, 'Tính năng nổi bật',
             font_size=44, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(sl, 2.0, 4.0, 9.3, 1.0,
             '10 nhóm tính năng chính — tất cả đã hoàn thiện, sẵn sàng sử dụng',
             font_size=18, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
add_slide_number(sl, 9)

# ═══════════════════════════════════════════════════════════
# SLIDE 10: TÍNH NĂNG CỐT LÕI
# ═══════════════════════════════════════════════════════════
features = [
    ('🤖', 'Trợ lý AI thông minh', 'Chat streaming, quiz inline,\ntóm tắt, giải thích gợi mở'),
    ('📚', 'Quản lý học tập', 'CRUD môn/khoá/bài học,\nupload & xem PDF'),
    ('📝', 'Bài tập & Chấm điểm', 'ĐÃ GIAO→ĐÃ NHẬN,\nAI chấm tự luận, audit log'),
    ('📊', 'Theo dõi tiến độ', 'Phiên học tự động, biểu đồ\n7 ngày, bảng xếp hạng'),
    ('🏆', 'Học mà chơi', 'Kim cương, streak,\n8 thành tựu, lịch sử GD'),
    ('🎯', 'Chẩn đoán điểm yếu', 'Phát hiện lỗ hổng, AI tạo\nbài tập cải thiện riêng'),
    ('🗺️', 'Mind Map & Graph', 'Sơ đồ tư duy AI, đồ thị\nkiến thức trực quan'),
    ('🃏', 'Flashcards (SM-2)', 'Thẻ ghi nhớ lặp lại\nngắt quãng, import Anki'),
]
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
add_text_box(sl, 1.0, 0.15, 11.3, 0.7, 'Nhóm tính năng cốt lõi',
             font_size=32, bold=True, align=PP_ALIGN.CENTER)
for i, (icon, title, desc) in enumerate(features):
    row = i // 4
    col = i % 4
    x = 0.4 + col * 3.25
    y = 1.1 + row * 3.1
    add_rect(sl, x, y, 3.0, 2.8, corner_radius=0.12)
    add_text_box(sl, x, y + 0.15, 3.0, 0.5, icon, font_size=26, align=PP_ALIGN.CENTER)
    add_text_box(sl, x, y + 0.8, 3.0, 0.4, title, font_size=14, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(sl, x + 0.15, y + 1.3, 2.7, 1.2, desc, font_size=10.5, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
add_slide_number(sl, 10)

# ═══════════════════════════════════════════════════════════
# SLIDE 11: AI TUTOR DEEP DIVE
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
add_text_box(sl, 0.5, 0.15, 12.3, 0.7, 'Trợ lý AI thông minh — Trái tim của EduSelf',
             font_size=32, bold=True, align=PP_ALIGN.CENTER)
# Flow
flow_steps = [
    ('📖', 'Đọc tài liệu', 'PDF Bài học'),
    ('💬', 'Chat với AI', 'trò chuyện mượt mà'),
    ('❓', 'Quiz tương tác', 'Inline trong chat'),
    ('📝', 'Bài tập tự luận', 'AI chấm điểm'),
    ('📊', 'Báo cáo & Cải thiện', 'Lộ trình cá nhân'),
]
for i, (icon, label, sub) in enumerate(flow_steps):
    x = 0.5 + i * 2.6
    add_rect(sl, x, 1.2, 2.3, 1.4, corner_radius=0.1)
    add_text_box(sl, x, 1.25, 2.3, 0.4, icon, font_size=24, align=PP_ALIGN.CENTER)
    add_text_box(sl, x, 1.7, 2.3, 0.3, label, font_size=12, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(sl, x, 2.1, 2.3, 0.3, sub, font_size=9, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    if i < 4:
        add_text_box(sl, x + 2.3, 1.5, 0.4, 0.4, '→', font_size=20, color=BLUE, align=PP_ALIGN.CENTER)
# Detail cards
details = [
    ('🧠', 'Phương pháp gợi mở', 'AI không đưa đáp án trực tiếp mà đặt câu hỏi gợi mở, giúp học sinh tự suy luận.'),
    ('📐', 'Hỗ trợ công thức toán học', 'Công thức toán, lý, hoá hiển thị sắc nét ngay trong chat.'),
    ('📋', 'Sinh quiz tự động (5-20 câu)', 'AI tạo câu hỏi từ nội dung bài học. Quiz kết thúc giúp đánh giá hiểu bài.'),
    ('🛡️', 'Kiểm soát ổn định', '30 req/phút chat, 10 req/phút generation. Bảo vệ hệ thống, trải nghiệm mượt.'),
]
for i, (icon, title, desc) in enumerate(details):
    row = i // 2
    col = i % 2
    x = 0.5 + col * 6.3
    y = 3.0 + row * 2.1
    add_rect(sl, x, y, 5.9, 1.8, corner_radius=0.12)
    add_text_box(sl, x + 0.15, y + 0.15, 0.5, 0.5, icon, font_size=20)
    add_text_box(sl, x + 0.75, y + 0.15, 5.0, 0.4, title, font_size=14, color=TEXT_WHITE, bold=True)
    add_text_box(sl, x + 0.75, y + 0.65, 5.0, 0.9, desc, font_size=11, color=TEXT_GRAY)
add_slide_number(sl, 11)

# ═══════════════════════════════════════════════════════════
# SLIDE 12: ASSIGNMENT WORKFLOW
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
add_text_box(sl, 0.5, 0.15, 12.3, 0.7, 'Quy trình Bài tập & Chấm điểm',
             font_size=32, bold=True, align=PP_ALIGN.CENTER)
wf_steps = [
    ('📝', 'ĐÃ GIAO', 'Giáo viên giao bài'),
    ('📤', 'ĐÃ NỘP', 'Học sinh nộp bài'),
    ('✅', 'ĐÃ CHẤM', 'AI / GV chấm điểm'),
    ('↩️', 'ĐÃ TRẢ', 'Trả bài cho HS'),
    ('🎯', 'ĐÃ NHẬN', 'HS chấp nhận điểm'),
]
for i, (icon, label, sub) in enumerate(wf_steps):
    x = 0.3 + i * 2.6
    add_rect(sl, x, 1.2, 2.3, 1.4, corner_radius=0.1)
    add_text_box(sl, x, 1.25, 2.3, 0.4, icon, font_size=24, align=PP_ALIGN.CENTER)
    add_text_box(sl, x, 1.7, 2.3, 0.3, label, font_size=12, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(sl, x, 2.1, 2.3, 0.3, sub, font_size=9, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    if i < 4:
        add_text_box(sl, x + 2.3, 1.55, 0.4, 0.4, '→', font_size=20, color=BLUE, align=PP_ALIGN.CENTER)
# Detail cards
assn_details = [
    ('🤖', 'Chấm điểm AI tự động', 'AI phân tích bài làm, đối chiếu rubric, đưa điểm + nhận xét. GV điều chỉnh trước khi trả bài.'),
    ('📋', 'Ma trận đề thi chuẩn MOET', 'Xây dựng ma trận đề theo chuẩn Bộ GD&ĐT. AI hỗ trợ tạo câu hỏi theo từng ô.'),
    ('📜', 'Audit Log minh bạch', 'Mọi thao tác đều được ghi nhận: ai giao, nộp, chấm, sửa điểm. Truy vết đầy đủ.'),
    ('📥', 'Xuất CSV bảng điểm', 'Xuất bảng điểm ra CSV để lưu trữ hoặc import vào phần mềm quản lý nhà trường.'),
]
for i, (icon, title, desc) in enumerate(assn_details):
    row = i // 2
    col = i % 2
    x = 0.5 + col * 6.3
    y = 3.0 + row * 2.1
    add_rect(sl, x, y, 5.9, 1.8, corner_radius=0.12)
    add_text_box(sl, x + 0.15, y + 0.15, 0.5, 0.5, icon, font_size=20)
    add_text_box(sl, x + 0.75, y + 0.15, 5.0, 0.4, title, font_size=14, color=TEXT_WHITE, bold=True)
    add_text_box(sl, x + 0.75, y + 0.65, 5.0, 0.9, desc, font_size=11, color=TEXT_GRAY)
add_slide_number(sl, 12)

# ═══════════════════════════════════════════════════════════
# SLIDE 13: GAMIFICATION
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
add_text_box(sl, 1.0, 0.15, 11.3, 0.7, 'Học mà chơi — Biến học tập thành cuộc phiêu lưu',
             font_size=30, bold=True, align=PP_ALIGN.CENTER)
g_stats = [
    ('💎', 'Kim cương thưởng\nkhi học, quiz, nộp bài'),
    ('🔥', 'Streak chuỗi ngày\nhọc liên tục'),
    ('🏅', '8 loại thành tựu\ncho các cột mốc'),
    ('🏆', 'Bảng xếp hạng\ntuần / tháng / tất cả'),
]
for i, (icon, label) in enumerate(g_stats):
    x = 1.0 + i * 3.1
    add_text_box(sl, x, 1.2, 2.8, 0.6, icon, font_size=38, align=PP_ALIGN.CENTER)
    add_text_box(sl, x, 1.9, 2.8, 0.8, label, font_size=12, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

g_details = [
    ('💰', 'Kiếm kim cương', 'Học bài, quiz, nộp bài\nđều được thưởng'),
    ('📋', 'Lịch sử giao dịch', 'Xem lại toàn bộ\nlịch sử kim cương'),
    ('⭐', 'Thành tựu', '8 thành tựu: học 10h,\nstreak 7 ngày, điểm cao...'),
    ('👑', 'Admin trao thưởng', 'Admin trao kim cương\ncho học sinh xuất sắc'),
]
for i, (icon, title, desc) in enumerate(g_details):
    x = 1.0 + i * 3.1
    y = 3.0
    add_rect(sl, x, y, 2.8, 2.5, corner_radius=0.12)
    add_text_box(sl, x, y + 0.15, 2.8, 0.5, icon, font_size=28, align=PP_ALIGN.CENTER)
    add_text_box(sl, x, y + 0.8, 2.8, 0.4, title, font_size=14, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(sl, x + 0.15, y + 1.3, 2.5, 1.0, desc, font_size=11, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
add_slide_number(sl, 13)

# ═══════════════════════════════════════════════════════════
# SLIDE 14: PROGRESS & WEAKNESS
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
add_text_box(sl, 0.5, 0.15, 12.3, 0.7, 'Theo dõi tiến độ & Chẩn đoán điểm yếu',
             font_size=28, bold=True, align=PP_ALIGN.CENTER)
# Left panel
add_rect(sl, 0.3, 1.2, 6.1, 5.2, corner_radius=0.15)
progress_items = [
    'Phiên học tự động (bắt đầu / tạm dừng / kết thúc)',
    'sendBeacon khi đóng tab — không mất dữ liệu',
    'Thống kê: tổng giờ, hôm nay, tuần này, điểm TB',
    'Biểu đồ cột 7 ngày học tập trực quan',
    'Bảng xếp hạng: tuần / tháng / tất cả',
    'Mục tiêu học tập: 2 giờ/ngày, 10 giờ/tuần',
]
add_text_box(sl, 0.5, 1.35, 5.5, 0.4, '📊 Theo dõi tiến độ', font_size=18, color=BLUE, bold=True)
for j, item in enumerate(progress_items):
    add_text_box(sl, 0.6, 1.95 + j * 0.7, 5.5, 0.5, f'✓  {item}', font_size=11.5, color=TEXT_GRAY)
# Right panel
add_rect(sl, 6.9, 1.2, 6.1, 5.2, corner_radius=0.15)
weakness_items = [
    'Tự động ghi nhận khi trả lời sai / điểm thấp',
    'Gom nhóm thông minh theo bài học + chủ đề',
    'Tự động cải thiện khi học sinh trả lời đúng',
    'AI tạo bài tập cải thiện riêng cho từng điểm yếu',
    'Giáo viên thêm ghi chú coach cho học sinh',
    'Dashboard trực quan — biết ngay cần ôn gì',
]
add_text_box(sl, 7.1, 1.35, 5.5, 0.4, '🎯 Chẩn đoán điểm yếu', font_size=18, color=PURPLE, bold=True)
for j, item in enumerate(weakness_items):
    add_text_box(sl, 7.2, 1.95 + j * 0.7, 5.5, 0.5, f'✓  {item}', font_size=11.5, color=TEXT_GRAY)
add_slide_number(sl, 14)

# ═══════════════════════════════════════════════════════════
# SLIDE 15: SECTION — Công nghệ
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_section_bg(sl, 7)
add_text_box(sl, 0, 2.8, 13.333, 1.0, 'Công nghệ sử dụng',
             font_size=44, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(sl, 2.0, 4.0, 9.3, 1.0,
             'Kiến trúc hiện đại, hiệu năng cao, dễ mở rộng — sẵn sàng cho quy mô toàn quốc',
             font_size=18, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
add_slide_number(sl, 15)

# ═══════════════════════════════════════════════════════════
# SLIDE 16: TECH STACK
# ═══════════════════════════════════════════════════════════
techs = [
    ('Trang web', 'Hiện đại, nhanh'), ('Máy chủ', 'Mạnh mẽ, ổn định'),
    ('Cơ sở dữ liệu', 'Lưu trữ an toàn'), ('Bảo mật', 'Đăng nhập & phân quyền'),
    ('Đám mây', 'Lưu trữ tài liệu'), ('Trí tuệ nhân tạo', 'AI thông minh'),
    ('Giao diện', 'Đẹp, dễ dùng'), ('Mọi thiết bị', 'Điện thoại, máy tính'),
]
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
add_text_box(sl, 1.0, 0.15, 11.3, 0.7, 'Kiến trúc công nghệ hiện đại',
             font_size=32, bold=True, align=PP_ALIGN.CENTER)
for i, (name, role) in enumerate(techs):
    row = i // 6
    col = i % 6
    x = 0.5 + col * 2.15
    y = 1.2 + row * 2.4
    add_rect(sl, x, y, 1.95, 2.0, corner_radius=0.1)
    add_text_box(sl, x, y + 0.3, 1.95, 0.6, name, font_size=13, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(sl, x, y + 1.0, 1.95, 0.4, role, font_size=9, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
# Extra tools
tools = 'Toàn bộ giao diện tiếng Việt • Dùng được trên mọi thiết bị • Hoạt động 24/7'
add_text_box(sl, 0.5, 6.4, 12.3, 0.4, tools, font_size=11, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
add_slide_number(sl, 16)

# ═══════════════════════════════════════════════════════════
# SLIDE 17: ARCHITECTURE
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
add_text_box(sl, 1.5, 0.15, 10.3, 0.7, 'Kiến trúc hệ thống',
             font_size=32, bold=True, align=PP_ALIGN.CENTER)
arch_boxes = [
    (1.5, 1.1, 10.3, 0.8, '🖥️📱  Client  —  Next.js 16 · React 19 · Tailwind CSS 4'),
    (1.5, 2.15, 10.3, 0.8, '▲  Vercel  —  Hosting Frontend + API Proxy'),
    (1.5, 3.4, 3.1, 1.0, '⚙️ Go Backend\nChi Router · GORM · JWT'),
    (5.1, 3.4, 3.1, 1.0, '🗄️ PostgreSQL 16\nSupabase Managed'),
    (8.7, 3.4, 3.1, 1.0, '☁️ Cloudflare R2\nLưu trữ PDF'),
    (1.5, 4.65, 10.3, 0.8, '🤖  Gemini API  —  AI Chat · Quiz · Chấm điểm · Báo cáo'),
]
for x, y, w, h, text in arch_boxes:
    add_rect(sl, x, y, w, h, corner_radius=0.1)
    add_text_box(sl, x + 0.15, y + 0.05, w - 0.3, h - 0.1, text, font_size=12, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
# Arrows
for arrow_y in [1.95, 3.2]:
    add_text_box(sl, 6.0, arrow_y, 1.3, 0.4, '⬇', font_size=20, color=BLUE, align=PP_ALIGN.CENTER)
add_slide_number(sl, 17)

# ═══════════════════════════════════════════════════════════
# SLIDE 18: SECTION — Điểm mới
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_section_bg(sl, 7)
add_text_box(sl, 0, 2.8, 13.333, 1.0, 'Điểm mới & Sáng tạo',
             font_size=44, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(sl, 2.0, 4.0, 9.3, 1.5,
             'Những yếu tố khác biệt cốt lõi làm nên giá trị riêng của EduSelf so với các nền tảng hiện có',
             font_size=18, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
add_slide_number(sl, 18)

# ═══════════════════════════════════════════════════════════
# SLIDE 19: COMPARISON TABLE
# ═══════════════════════════════════════════════════════════
cmp_data = [
    ('AI Tutor', 'Chatbot trả lời trực tiếp,\nkhông có phương pháp sư phạm', 'Trợ lý AI thông minh — đặt câu hỏi\ngợi mở, dẫn dắt tự khám phá'),
    ('Cá nhân hoá', 'Gợi ý khoá học dựa trên\ndanh mục có sẵn', 'AI phân tích năng lực → chẩn đoán\nđiểm yếu → bài tập riêng biệt'),
    ('Chấm điểm', 'Trắc nghiệm tự động,\ntự luận chấm thủ công', 'AI chấm cả tự luận theo rubric,\ngiáo viên chỉ kiểm tra lại'),
    ('Học mà chơi', 'Huy hiệu cơ bản,\nkhông có hệ thống thưởng', 'Kim cương + Streak + 8 Thành tựu\n+ Bảng xếp hạng đầy đủ'),
    ('Công cụ học', 'Video + quiz là chính', 'Mind Map AI + Knowledge Graph +\nFlashcards (SM-2) + Sơ đồ tư duy'),
    ('Báo cáo', 'Bảng điểm tĩnh,\nkhông phân tích', 'AI tạo báo cáo hàng tuần +\nDashboard trực quan theo vai trò'),
    ('Ngôn ngữ', 'Chủ yếu tiếng Anh,\nít hỗ trợ tiếng Việt', 'Toàn bộ UI bằng tiếng Việt,\nphù hợp chương trình Việt Nam'),
    ('Phân quyền', '2-3 vai trò cơ bản', 'phân quyền 5 cấp + ma trận\nphân quyền chi tiết'),
]
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
add_text_box(sl, 0.5, 0.1, 12.3, 0.6, 'Điều gì làm EduSelf khác biệt?',
             font_size=30, bold=True, align=PP_ALIGN.CENTER)
# Table header
add_rect(sl, 0.3, 0.9, 2.5, 0.6, fill_color=RGBColor(0x0D, 0x18, 0x3A))
add_text_box(sl, 0.3, 0.95, 2.5, 0.5, 'Tiêu chí', font_size=13, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_rect(sl, 2.85, 0.9, 5.0, 0.6, fill_color=RGBColor(0x0D, 0x18, 0x3A))
add_text_box(sl, 2.85, 0.95, 5.0, 0.5, 'Nền tảng khác', font_size=13, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_rect(sl, 7.9, 0.9, 5.1, 0.6, fill_color=RGBColor(0x0D, 0x18, 0x3A))
add_text_box(sl, 7.9, 0.95, 5.1, 0.5, 'EduSelf', font_size=13, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
for i, (criterion, old, new) in enumerate(cmp_data):
    y = 1.55 + i * 0.72
    bg_color = RGBColor(0x0A, 0x0F, 0x20) if i % 2 == 0 else RGBColor(0x0E, 0x13, 0x28)
    add_rect(sl, 0.3, y, 2.5, 0.68, fill_color=bg_color)
    add_text_box(sl, 0.4, y + 0.1, 2.3, 0.5, criterion, font_size=11, color=TEXT_WHITE, bold=True)
    add_rect(sl, 2.85, y, 5.0, 0.68, fill_color=bg_color)
    add_text_box(sl, 3.0, y + 0.05, 4.7, 0.6, old, font_size=10, color=RED_LIGHT)
    add_rect(sl, 7.9, y, 5.1, 0.68, fill_color=bg_color)
    add_text_box(sl, 8.05, y + 0.05, 4.8, 0.6, new, font_size=10, color=GREEN_LIGHT)
add_slide_number(sl, 19)

# ═══════════════════════════════════════════════════════════
# SLIDE 20: HIỆU QUẢ
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
add_text_box(sl, 1.5, 0.15, 10.3, 0.7, 'Hiệu quả mang lại',
             font_size=34, bold=True, align=PP_ALIGN.CENTER)
impacts = [
    ('⚙️', 'Hiệu quả Kỹ thuật', BLUE, [
        'Tiết kiệm 70% thời gian soạn & chấm bài',
        'AI 24/7, phản hồi tức thì qua SSE',
        '45+ API endpoints, kiến trúc monorepo',
        'sendBeacon — không mất dữ liệu',
    ]),
    ('💰', 'Hiệu quả Kinh tế', GREEN, [
        'Chi phí thấp hơn giải pháp nước ngoài',
        'Giảm chi phí in ấn đề thi, tài liệu',
        'Mở rộng không giới hạn nhờ cloud',
        'Không cần mua bản quyền đắt đỏ',
    ]),
    ('🌍', 'Hiệu quả Xã hội', PURPLE, [
        'Thu hẹp khoảng cách giáo dục',
        'Học sinh vùng sâu tiếp cận AI Tutor',
        'Tăng cường sự tham gia của phụ huynh',
        'Giảm áp lực công việc cho giáo viên',
    ]),
]
for i, (icon, title, clr, items) in enumerate(impacts):
    x = 0.5 + i * 4.25
    add_rect(sl, x, 1.1, 3.95, 5.0, corner_radius=0.15)
    add_text_box(sl, x, 1.2, 3.95, 0.5, icon, font_size=32, align=PP_ALIGN.CENTER)
    add_text_box(sl, x, 1.8, 3.95, 0.4, title, font_size=17, color=clr, bold=True, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        add_text_box(sl, x + 0.2, 2.4 + j * 0.8, 3.55, 0.7, f'▸  {item}', font_size=11, color=TEXT_GRAY)
add_slide_number(sl, 20)

# ═══════════════════════════════════════════════════════════
# SLIDE 21: KHẢ NĂNG ÁP DỤNG
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
add_text_box(sl, 1.0, 0.15, 11.3, 0.7, 'Khả năng áp dụng',
             font_size=34, bold=True, align=PP_ALIGN.CENTER)
# Left
add_rect(sl, 0.3, 1.2, 6.1, 3.5, corner_radius=0.15)
add_text_box(sl, 0.5, 1.35, 5.7, 0.4, '🎯 Đối tượng sử dụng', font_size=17, color=BLUE, bold=True)
targets = [
    '🏫  Trường học các cấp (Tiểu học, THCS, THPT)',
    '🏠  Học sinh tự học tại nhà',
    '👨‍🏫  Trung tâm dạy thêm, gia sư',
    '🏢  Sở / Phòng GD&ĐT quản lý tập trung',
]
for j, t in enumerate(targets):
    add_text_box(sl, 0.6, 2.0 + j * 0.65, 5.5, 0.5, t, font_size=13, color=TEXT_GRAY)
# Right
add_rect(sl, 6.9, 1.2, 6.1, 3.5, corner_radius=0.15)
add_text_box(sl, 7.1, 1.35, 5.7, 0.4, '📈 Quy mô triển khai', font_size=17, color=GREEN, bold=True)
scales = [
    '🏫  1 trường — Triển khai trong 1-2 ngày',
    '🏘️  Cụm trường — Dùng chung hạ tầng cloud',
    '🏙️  Cấp tỉnh — Azure Container Apps auto-scale',
    '🌏  Toàn quốc — Kiến trúc sẵn sàng mở rộng',
]
for j, s in enumerate(scales):
    add_text_box(sl, 7.2, 2.0 + j * 0.65, 5.5, 0.5, s, font_size=13, color=TEXT_GRAY)
# Future
add_rect(sl, 0.3, 5.2, 12.7, 1.3, fill_color=RGBColor(0x15, 0x10, 0x05),
         border_color=RGBColor(0x40, 0x30, 0x10), corner_radius=0.1)
add_text_box(sl, 0.5, 5.3, 0.6, 0.6, '🔮', font_size=30, align=PP_ALIGN.CENTER)
add_text_box(sl, 1.1, 5.3, 11.5, 0.4, 'Định hướng phát triển', font_size=14, color=AMBER, bold=True)
add_text_box(sl, 1.1, 5.75, 11.5, 0.5,
             'Voice Chat AI • Quiz Battle • Study Planner • Tích hợp SIS • Mobile App iOS/Android',
             font_size=12, color=TEXT_GRAY)
add_slide_number(sl, 21)

# ═══════════════════════════════════════════════════════════

# ═══════════════════════════════════════════════════════════
# SLIDE 22: SECTION — Hành trình phát triển
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_section_bg(sl, 8)
add_text_box(sl, 0, 2.8, 13.333, 1.0, 'Hành trình phát triển',
             font_size=44, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(sl, 1.5, 4.0, 10.3, 1.5,
             'Từ ý tưởng ban đầu đến sản phẩm hoàn thiện trong gần 6 tuần làm việc liên tục',
             font_size=18, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
add_slide_number(sl, 22)

# ═══════════════════════════════════════════════════════════
# SLIDE 23: TIMELINE
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
add_text_box(sl, 0.5, 0.15, 12.3, 0.7, 'Từ ý tưởng đến sản phẩm',
             font_size=32, bold=True, align=PP_ALIGN.CENTER)
timeline = [
    ('1', BLUE, 'Tuần 1 (7/5 – 14/5/2026) — Khởi động & Nền móng',
     'Khởi tạo Next.js + Go. Xác thực, phân quyền, trình xem PDF, AI Tutor nền tảng.'),
    ('2', RGBColor(0x63,0x66,0xF1), 'Tuần 2 (15/5 – 21/5/2026) — AI & Học tập thông minh',
     'AI Socratic Tutor, streaming chat, quiz, chấm điểm. Flashcards, Mind Map, Knowledge Graph.'),
    ('3', PURPLE, 'Tuần 3 (22/5 – 28/5/2026) — Hoàn thiện hệ sinh thái',
     'Module Phụ huynh, Kế hoạch học tập, Ngân hàng câu hỏi, Điểm danh, Báo cáo AI.'),
    ('4', GREEN, 'Tuần 4-5 (29/5 – 15/6) — Triển khai & Hoàn thiện',
     'Deploy Vercel + Azure. Kiểm thử, sửa lỗi, tối ưu. Hoàn thiện UI/UX tiếng Việt.'),
]
for i, (num, clr, title, desc) in enumerate(timeline):
    y = 1.1 + i * 1.45
    add_rect(sl, 0.5, y, 12.3, 1.2, corner_radius=0.1)
    add_circle(sl, 0.7, y + 0.35, 0.5, clr)
    add_text_box(sl, 0.7, y + 0.35, 0.5, 0.5, num, font_size=16, color=RGBColor(0xFF,0xFF,0xFF), bold=True, align=PP_ALIGN.CENTER)
    add_text_box(sl, 1.45, y + 0.1, 11.0, 0.35, title, font_size=14, color=TEXT_WHITE, bold=True)
    add_text_box(sl, 1.45, y + 0.55, 11.0, 0.5, desc, font_size=11, color=TEXT_GRAY)
# Launch item
y_launch = 1.1 + len(timeline) * 1.45 + 0.1
add_rect(sl, 0.5, y_launch, 12.3, 1.0, fill_color=RGBColor(0x15,0x10,0x05),
         border_color=RGBColor(0x40,0x30,0x10), corner_radius=0.1)
add_text_box(sl, 0.7, y_launch + 0.15, 11.5, 0.35, '🚀  Hiện tại — Sẵn sàng ra mắt', font_size=14, color=AMBER, bold=True)
add_text_box(sl, 0.7, y_launch + 0.5, 11.5, 0.35, '10 nhóm tính năng, 45+ dịch vụ, 5 vai trò. Sẵn sàng triển khai cho trường học.', font_size=11, color=TEXT_GRAY)
add_slide_number(sl, 23)

# ═══════════════════════════════════════════════════════════
# SLIDE 24: QUY TRÌNH SỬ DỤNG
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
add_text_box(sl, 0.5, 0.15, 12.3, 0.7, 'Một buổi học điển hình trên EduSelf',
             font_size=30, bold=True, align=PP_ALIGN.CENTER)
flow_items = [
    ('🔑', 'Đăng nhập', 'Mở máy, vào web'),
    ('📊', 'Dashboard', 'Xem tiến độ'),
    ('📖', 'Chọn bài', 'Môn → Khoá → Bài'),
    ('📄', 'Đọc tài liệu', 'Xem PDF'),
    ('💬', 'Chat với AI', 'Hỏi đáp gợi mở'),
    ('✅', 'Làm quiz', 'Kiểm tra hiểu bài'),
    ('💎', 'Nhận thưởng', 'Kim cương & streak'),
]
for i, (icon, label, sub) in enumerate(flow_items):
    x = 0.3 + i * 1.85
    add_rect(sl, x, 1.1, 1.65, 1.3, corner_radius=0.1)
    add_text_box(sl, x, 1.15, 1.65, 0.4, icon, font_size=22, align=PP_ALIGN.CENTER)
    add_text_box(sl, x, 1.6, 1.65, 0.3, label, font_size=11, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(sl, x, 1.95, 1.65, 0.3, sub, font_size=9, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    if i < 6:
        add_text_box(sl, x + 1.65, 1.4, 0.25, 0.3, '→', font_size=16, color=BLUE, align=PP_ALIGN.CENTER)
# Left panel
add_rect(sl, 0.3, 2.8, 6.1, 3.8, corner_radius=0.15)
add_text_box(sl, 0.5, 2.95, 5.5, 0.35, '⏱️  Thời lượng: ~45 phút / bài', font_size=16, color=BLUE, bold=True)
time_items = [
    '15 phút đầu: Đọc tài liệu PDF, đánh dấu',
    '15 phút giữa: Trò chuyện với AI — đặt câu hỏi',
    '10 phút cuối: Làm quiz kết thúc bài',
    '5 phút: Nhận kim cương, kiểm tra streak',
]
for j, item in enumerate(time_items):
    add_text_box(sl, 0.6, 3.45 + j * 0.7, 5.5, 0.5, f'✓  {item}', font_size=12, color=TEXT_GRAY)
# Right panel
add_rect(sl, 6.9, 2.8, 6.1, 3.8, corner_radius=0.15)
add_text_box(sl, 7.1, 2.95, 5.5, 0.35, '✨  Điểm đặc biệt', font_size=16, color=GREEN, bold=True)
special_items = [
    'AI không trả lời thẳng — gợi mở để tự nghĩ',
    'Công thức toán hiển thị đẹp trong chat',
    'Quiz hiện ngay trong khung chat',
    'Sai → AI ghi nhận điểm yếu, tạo bài tập riêng',
]
for j, item in enumerate(special_items):
    add_text_box(sl, 7.2, 3.45 + j * 0.7, 5.5, 0.5, f'✓  {item}', font_size=12, color=TEXT_GRAY)
add_slide_number(sl, 24)

# ═══════════════════════════════════════════════════════════
# SLIDE 25: CHI PHÍ TRIỂN KHAI
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
add_text_box(sl, 1.0, 0.15, 11.3, 0.7, 'Chi phí triển khai',
             font_size=34, bold=True, align=PP_ALIGN.CENTER)
cost_data = [
    ('Chi phí nền tảng / năm', '~3-5 triệu VNĐ (cloud + domain)', '~20-50 triệu (bản quyền + hosting)'),
    ('Chi phí AI / năm', '~2-3 triệu VNĐ (API)', 'Không có sẵn (mua thêm)'),
    ('Triển khai', '1-2 ngày, miễn phí', '1-4 tuần, cần chuyên gia'),
    ('Bảo trì', 'Tự động qua cloud', 'Thủ công / thuê ngoài'),
    ('Ngôn ngữ', 'Tiếng Việt 100%', 'Tiếng Anh là chính'),
    ('Tổng chi phí / năm', '5-8 triệu VNĐ', '30-80 triệu VNĐ'),
]
# Header
add_rect(sl, 0.3, 1.1, 3.5, 0.6, fill_color=RGBColor(0x0D,0x18,0x3A))
add_text_box(sl, 0.3, 1.15, 3.5, 0.5, 'Hạng mục', font_size=13, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_rect(sl, 3.85, 1.1, 4.5, 0.6, fill_color=RGBColor(0x0D,0x18,0x3A))
add_text_box(sl, 3.85, 1.15, 4.5, 0.5, 'EduSelf', font_size=13, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_rect(sl, 8.4, 1.1, 4.6, 0.6, fill_color=RGBColor(0x0D,0x18,0x3A))
add_text_box(sl, 8.4, 1.15, 4.6, 0.5, 'Giải pháp khác', font_size=13, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
for i, (item, es, other) in enumerate(cost_data):
    y = 1.75 + i * 0.75
    bg = RGBColor(0x0A,0x0F,0x20) if i % 2 == 0 else RGBColor(0x0E,0x13,0x28)
    add_rect(sl, 0.3, y, 3.5, 0.7, fill_color=bg)
    add_text_box(sl, 0.4, y + 0.1, 3.3, 0.5, item, font_size=12, color=TEXT_WHITE, bold=True)
    add_rect(sl, 3.85, y, 4.5, 0.7, fill_color=bg)
    add_text_box(sl, 4.0, y + 0.1, 4.2, 0.5, es, font_size=11, color=GREEN_LIGHT)
    add_rect(sl, 8.4, y, 4.6, 0.7, fill_color=bg)
    add_text_box(sl, 8.55, y + 0.1, 4.3, 0.5, other, font_size=11, color=RED_LIGHT)
# Note
add_text_box(sl, 0.5, 6.5, 12.3, 0.6,
             'Ước tính cho quy mô 1 trường (~500-1000 học sinh). EduSelf tiết kiệm 80-90% chi phí nhờ hạ tầng đám mây và công nghệ nguồn mở.',
             font_size=12, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
add_slide_number(sl, 25)

# SLIDE 22: THANK YOU
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
add_text_box(sl, 0, 1.8, 13.333, 1.0, '🙏', font_size=60, align=PP_ALIGN.CENTER)
add_text_box(sl, 0, 3.2, 13.333, 1.2, 'Xin cảm ơn!',
             font_size=58, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(sl, 0, 4.6, 13.333, 0.8,
             'Hệ thống Quản lý Học tập Thông minh\nSẵn sàng cho tương lai giáo dục Việt Nam',
             font_size=20, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
for i, (text, clr) in enumerate([
    ('🤖 AI-Powered', RGBColor(0xC4, 0xB5, 0xFD)),
    ('🇻🇳 Made in Vietnam', RGBColor(0x93, 0xC5, 0xFD)),
    ('📚 Education First', RGBColor(0x6E, 0xE7, 0xB7)),
]):
    add_text_box(sl, 3.5 + i * 2.5, 5.8, 2.3, 0.4, text, font_size=12, color=clr, align=PP_ALIGN.CENTER)
add_slide_number(sl, 22)

# ── Save ──
output_path = '/Users/hoangnamtran/WORKS/lms/presentation.pptx'
prs.save(output_path)
print(f'PPTX saved to {output_path}')
print(f'Total slides: {len(prs.slides)}')
